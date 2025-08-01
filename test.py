from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation and slide
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout

# Set title
title = slide.shapes.title
title.text = "Work-Flow الگوریتم Q-Learning خروجی‌های اضطراری"

# Define bullet points for the workflow
content = slide.shapes.placeholders[1].text_frame
content.clear()

steps = [
    "1. تولید حالت اولیه (Initial State):\n   • تولید تصادفی پیکربندی k خروج با عرض ω بدون هم‌پوشانی",
    "2. انتخاب عمل (Action Selection):\n   • e-greedy: با احتمال ε عمل تصادفی، و با احتمال 1-ε بهترین عمل طبق Q",
    "3. اعمال عمل و ارزیابی (Apply & Evaluate):\n   • اجرای عمل، محاسبه هزینه جدید و پاداش = هزینه_قدیم - هزینه_جدید",
    "4. به‌روزرسانی جدول Q (Q-Update):\n   • استفاده از معادله بلمن: Q(s,a) ← Q + α[r + γ max Q(s',·) - Q]",
    "5. ثبت بهترین جواب و کاهش ε (Best & Decay):\n   • به‌روزرسانی best_solution و decay ε برای بهینه‌سازی بیشتر",
]

for step in steps:
    p = content.add_paragraph()
    p.text = step
    p.level = 0
    p.font.size = Pt(16)

# Save the presentation
pptx_path = "qlearning_exit_workflow.pptx"
prs.save(pptx_path)

print(f"[Download the slide here: {pptx_path}]")
